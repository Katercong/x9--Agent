"""X9 Database API server.

Run with:
    cd F:\\Claude_Project\\Database
    python -m uvicorn app.main:app --host 0.0.0.0 --port 18765 --reload

Then open http://localhost:18765/ in a browser.
(Port 8765 conflicted with another local project; default switched to 18765.)
"""
from __future__ import annotations
import json
import os
import sqlite3
from pathlib import Path
from typing import Any
from urllib.parse import unquote

import psycopg
from fastapi import FastAPI, HTTPException, Request, Response
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from psycopg.rows import dict_row

from app.v1 import router as v1_router
from app.llm import router as llm_router
from app.auth_routes import router as auth_router
from app.agent import router as agent_router
from app.outreach_ai import router as outreach_ai_router
from app.title_optimizer import router as title_router
from app.keyword_ai import router as keyword_ai_router

ROOT = Path(__file__).resolve().parent.parent
DB_PATH = ROOT / "database.db"
ASSETS_DIR = ROOT / "assets"
STATIC_DIR = ROOT / "app" / "static"
INTERN_A = Path(r"F:\实习生\A社媒")
PG_DSN = os.environ.get(
    "X9_PG_DSN",
    "postgresql://x9:x9_local_dev_2026@localhost:15432/x9db?connect_timeout=5",
)

# JSON columns: stored as JSON-encoded TEXT, surfaced as Python lists/dicts in API
PRODUCT_JSON_COLS = {
    "selling_points_en", "selling_points_zh", "pain_points_zh",
    "scenarios_en", "scenarios_zh", "vocabulary_en",
    "creative_angles_en", "safe_scenes_en", "creator_match_levels",
}
CREATOR_JSON_COLS = {"category_tags"}

# editable columns (whitelist) — limits accidental writes via PUT
PRODUCT_EDITABLE = {
    "name_en", "name_zh", "subcategory", "series", "size_label",
    "pcs_per_pack", "packs_per_case",
    "price_tiktok", "price_temu", "price_ebay", "price_ebay_local", "price_independent",
    "currency", "positioning_zh", "tier",
    "description_en", "description_zh",
    "selling_points_en", "selling_points_zh", "pain_points_zh",
    "scenarios_en", "scenarios_zh",
    "target_audience_en", "target_audience_zh",
    "proof", "vocabulary_en", "creative_angles_en", "safe_scenes_en",
    "focus_zh", "amazon_url", "short_url", "tk_content_key",
    "commission_rate_default", "creator_match_levels",
    "creator_persona_zh", "is_main_push", "status",
}
CREATOR_EDITABLE = {
    "platform", "profile_url", "display_name", "country", "language",
    "category_tags", "followers", "followers_raw", "tier",
    "avg_views", "gmv_30d_usd", "pps", "sample_score", "post_rate_est",
    "email", "whatsapp", "instagram_handle", "youtube_handle",
    "current_status", "store_assigned", "owner_bd",
    "first_contact_date", "last_contact_date", "notes", "source",
    "quality_score",
}
OUTREACH_EDITABLE = {
    "creator_id", "event_date", "store_name", "bd_owner",
    "action", "status", "channel", "message",
    "sample_qty", "commission_rate", "video_url", "ad_auth_code", "remark",
}


def get_con() -> sqlite3.Connection:
    con = sqlite3.connect(DB_PATH)
    con.execute("PRAGMA foreign_keys=ON")
    con.row_factory = sqlite3.Row
    return con


def get_pg_con() -> psycopg.Connection:
    return psycopg.connect(PG_DSN, row_factory=dict_row)


def row_to_dict(r: sqlite3.Row, json_cols: set[str]) -> dict:
    d = dict(r)
    for k in json_cols:
        if k in d and isinstance(d[k], str) and d[k]:
            try:
                d[k] = json.loads(d[k])
            except json.JSONDecodeError:
                pass
        elif k in d and not d[k]:
            d[k] = []
    return d


def encode_for_db(payload: dict, json_cols: set[str]) -> dict:
    out = {}
    for k, v in payload.items():
        if k in json_cols and not isinstance(v, str):
            out[k] = json.dumps(v or [], ensure_ascii=False)
        else:
            out[k] = v
    return out


# ============================================================
# App
# ============================================================
app = FastAPI(title="X9 Database", version="3.7.0")
app.include_router(v1_router)
app.include_router(llm_router)
app.include_router(auth_router)
app.include_router(agent_router)
app.include_router(outreach_ai_router)
app.include_router(title_router)
app.include_router(keyword_ai_router)


@app.get("/")
def root() -> FileResponse:
    return FileResponse(STATIC_DIR / "index.html")


# ---------- New React preview UI (3 roles x 24 pages) ----------
# Sub-path deploy at /web-preview/. Old UI at /  is untouched.
# Build with: cd web && npm run build:deploy && npm run deploy
WEB_PREVIEW_DIR = STATIC_DIR / "web-preview"


@app.get("/web-preview")
def web_preview_index_redirect() -> FileResponse:
    return FileResponse(WEB_PREVIEW_DIR / "index.html")


@app.get("/web-preview/")
def web_preview_index() -> FileResponse:
    return FileResponse(WEB_PREVIEW_DIR / "index.html")


@app.get("/web-preview/{full_path:path}")
def web_preview_spa(full_path: str) -> FileResponse:
    """SPA fallback: serve actual file if exists, otherwise index.html so React Router can handle the path."""
    candidate = WEB_PREVIEW_DIR / full_path
    if candidate.is_file():
        return FileResponse(candidate)
    return FileResponse(WEB_PREVIEW_DIR / "index.html")


@app.get("/api/business-summary")
def business_summary() -> dict:
    try:
        from app.pg_dashboard import api_business_summary

        return api_business_summary()
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"business summary unavailable: {exc}") from exc


@app.get("/api/stats")
def stats() -> dict:
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute(
            """
            SELECT
              (SELECT COUNT(*)::int FROM product) AS products_total,
              (SELECT COUNT(*)::int FROM product WHERE COALESCE(is_main_push, 0) = 1) AS products_main_push,
              (SELECT COUNT(*)::int FROM creator) AS creators_total,
              (SELECT COUNT(*)::int FROM outreach) AS outreach_total
            """
        )
        out = cur.fetchone()
        cur.execute(
            """
            SELECT c.code, c.name_zh AS name, COUNT(p.id)::int AS count
            FROM category c
            LEFT JOIN product p ON p.category_id = c.id
            GROUP BY c.id, c.code, c.name_zh, c.sort_order
            ORDER BY c.sort_order NULLS LAST, c.id
            """
        )
        out["by_category"] = cur.fetchall()
        cur.execute(
            """
            SELECT COALESCE(current_status, '(unknown)') AS status, COUNT(*)::int AS count
            FROM creator
            GROUP BY current_status
            ORDER BY count DESC
            """
        )
        out["creators_by_status"] = cur.fetchall()
        cur.execute(
            """
            SELECT COALESCE(tier, '(unknown)') AS tier, COUNT(*)::int AS count
            FROM creator
            GROUP BY tier
            ORDER BY tier
            """
        )
        out["creators_by_tier"] = cur.fetchall()
        return out


@app.get("/api/categories")
def list_categories() -> list[dict]:
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute("SELECT * FROM category ORDER BY sort_order NULLS LAST, id")
        return cur.fetchall()


# ---------- Products ----------
@app.get("/api/products")
def list_products(category: str | None = None, q: str | None = None,
                  is_main_push: int | None = None, limit: int = 200, offset: int = 0) -> dict:
    where = ["1=1"]
    args: list[Any] = []
    if category:
        where.append("c.code = %s")
        args.append(category)
    if q:
        where.append("(p.sku_code ILIKE %s OR p.name_en ILIKE %s OR p.name_zh ILIKE %s OR p.series ILIKE %s)")
        like = f"%{q}%"
        args += [like, like, like, like]
    if is_main_push is not None:
        where.append("COALESCE(p.is_main_push, 0) = %s")
        args.append(int(is_main_push))
    where_sql = " AND ".join(where)
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute(
            f"""
            SELECT p.*, c.code AS category_code, c.name_zh AS category_name
            FROM product p
            LEFT JOIN category c ON p.category_id = c.id
            WHERE {where_sql}
            ORDER BY COALESCE(p.is_main_push, 0) DESC, p.id ASC
            LIMIT %s OFFSET %s
            """,
            args + [limit, offset],
        )
        rows = [row_to_dict(r, PRODUCT_JSON_COLS) for r in cur.fetchall()]
        cur.execute(
            f"""
            SELECT COUNT(*)::int AS total
            FROM product p
            LEFT JOIN category c ON p.category_id = c.id
            WHERE {where_sql}
            """,
            args,
        )
        total = cur.fetchone()["total"]
        if rows:
            ids = [r["id"] for r in rows]
            cur.execute(
                """
                SELECT product_id, rel_path, kind, caption, display_order
                FROM product_image
                WHERE product_id = ANY(%s)
                ORDER BY display_order NULLS LAST, id
                LIMIT 8000
                """,
                (ids,),
            )
            by_pid: dict[int, list] = {}
            for img in cur.fetchall():
                by_pid.setdefault(img["product_id"], []).append({
                    "url": rel_path_to_url(img["rel_path"] or ""),
                    "kind": img["kind"],
                    "caption": img["caption"],
                    "order": img["display_order"],
                })
            for r in rows:
                r["images"] = by_pid.get(r["id"], [])[:30]
        return {"total": total, "items": rows}


@app.get("/api/products/{sku}")
def get_product(sku: str) -> dict:
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute(
            """
            SELECT p.*, c.code AS category_code, c.name_zh AS category_name
            FROM product p
            LEFT JOIN category c ON p.category_id = c.id
            WHERE p.sku_code = %s
            """,
            (sku,),
        )
        r = cur.fetchone()
        if not r:
            raise HTTPException(404, "product not found")
        out = row_to_dict(r, PRODUCT_JSON_COLS)
        cur.execute(
            """
            SELECT id, rel_path, kind, caption, display_order
            FROM product_image
            WHERE product_id = %s
            ORDER BY display_order NULLS LAST, id
            """,
            (out["id"],),
        )
        out["images"] = [
            {
                "id": i["id"],
                "url": rel_path_to_url(i["rel_path"] or ""),
                "kind": i["kind"],
                "caption": i["caption"],
                "order": i["display_order"],
            }
            for i in cur.fetchall()
        ]
        return out


@app.put("/api/products/{sku}")
async def update_product(sku: str, request: Request) -> dict:
    payload = await request.json()
    payload = {k: v for k, v in payload.items() if k in PRODUCT_EDITABLE}
    if not payload:
        raise HTTPException(400, "no editable fields in payload")
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute("SELECT id FROM product WHERE sku_code = %s", (sku,))
        r = cur.fetchone()
        if not r:
            raise HTTPException(404, "product not found")
        payload = encode_for_db(payload, PRODUCT_JSON_COLS)
        sets = ", ".join([f"{k} = %s" for k in payload])
        cur.execute(
            f"UPDATE product SET {sets} WHERE sku_code = %s",
            list(payload.values()) + [sku],
        )
        cur.execute(
            """
            INSERT INTO audit_log(table_name, record_id, action, changes, operator)
            VALUES('product', %s, 'update', %s, %s)
            """,
            (r["id"], json.dumps(payload, ensure_ascii=False), request.headers.get("x-user", "")),
        )
    return get_product(sku)


@app.post("/api/products")
async def create_product(request: Request) -> dict:
    payload = await request.json()
    sku = payload.get("sku_code")
    if not sku:
        raise HTTPException(400, "sku_code required")
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute("SELECT 1 FROM product WHERE sku_code = %s", (sku,))
        if cur.fetchone():
            raise HTTPException(409, "sku_code already exists")
        cat = payload.get("category_code")
        if cat:
            cur.execute("SELECT id FROM category WHERE code = %s", (cat,))
            cat_row = cur.fetchone()
            if cat_row:
                payload["category_id"] = cat_row["id"]
        payload.pop("category_code", None)
        fields = {k: v for k, v in payload.items() if k in PRODUCT_EDITABLE | {"sku_code", "category_id"}}
        fields = encode_for_db(fields, PRODUCT_JSON_COLS)
        cols = list(fields.keys())
        placeholders = ",".join(["%s"] * len(cols))
        cur.execute(
            f"INSERT INTO product({','.join(cols)}) VALUES({placeholders}) RETURNING id",
            list(fields.values()),
        )
        pid = cur.fetchone()["id"]
        cur.execute(
            """
            INSERT INTO audit_log(table_name, record_id, action, changes, operator)
            VALUES('product', %s, 'insert', %s, %s)
            """,
            (pid, json.dumps(fields, ensure_ascii=False), request.headers.get("x-user", "")),
        )
    return get_product(sku)


# ---------- Creators ----------
@app.get("/api/creators")
def list_creators(q: str | None = None, tier: str | None = None,
                  status: str | None = None, owner: str | None = None,
                  limit: int = 200, offset: int = 0) -> dict:
    where = ["1=1"]
    args: list[Any] = []
    if q:
        where.append("(handle ILIKE %s OR display_name ILIKE %s OR notes ILIKE %s)")
        like = f"%{q}%"
        args += [like, like, like]
    if tier:
        where.append("tier = %s")
        args.append(tier)
    if status:
        where.append("current_status = %s")
        args.append(status)
    if owner:
        where.append("owner_bd = %s")
        args.append(owner)
    where_sql = " AND ".join(where)
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute(
            f"""
            SELECT *
            FROM creator
            WHERE {where_sql}
            ORDER BY updated_at DESC NULLS LAST, id DESC
            LIMIT %s OFFSET %s
            """,
            args + [limit, offset],
        )
        rows = [row_to_dict(r, CREATOR_JSON_COLS) for r in cur.fetchall()]
        cur.execute(f"SELECT COUNT(*)::int AS total FROM creator WHERE {where_sql}", args)
        total = cur.fetchone()["total"]
        return {"total": total, "items": rows}


@app.get("/api/creators/{cid}")
def get_creator(cid: int) -> dict:
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute("SELECT * FROM creator WHERE id = %s", (cid,))
        r = cur.fetchone()
        if not r:
            raise HTTPException(404, "creator not found")
        out = row_to_dict(r, CREATOR_JSON_COLS)
        cur.execute(
            """
            SELECT id, event_date, store_name, bd_owner, action, status, channel,
                   sample_qty, commission_rate, video_url, ad_auth_code, remark, created_at
            FROM outreach
            WHERE creator_id = %s
            ORDER BY event_date DESC NULLS LAST, id DESC
            """,
            (cid,),
        )
        out["outreach"] = cur.fetchall()
        cur.execute(
            """
            SELECT cp.product_id, p.sku_code, p.name_en, cp.relation, cp.note
            FROM creator_product cp
            JOIN product p ON p.id = cp.product_id
            WHERE cp.creator_id = %s
            """,
            (cid,),
        )
        out["products"] = cur.fetchall()
        return out


@app.put("/api/creators/{cid}")
async def update_creator(cid: int, request: Request) -> dict:
    payload = await request.json()
    payload = {k: v for k, v in payload.items() if k in CREATOR_EDITABLE}
    if not payload:
        raise HTTPException(400, "no editable fields")
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute("SELECT 1 FROM creator WHERE id = %s", (cid,))
        if not cur.fetchone():
            raise HTTPException(404, "creator not found")
        payload = encode_for_db(payload, CREATOR_JSON_COLS)
        sets = ", ".join([f"{k} = %s" for k in payload])
        cur.execute(f"UPDATE creator SET {sets} WHERE id = %s", list(payload.values()) + [cid])
        cur.execute(
            """
            INSERT INTO audit_log(table_name, record_id, action, changes, operator)
            VALUES('creator', %s, 'update', %s, %s)
            """,
            (cid, json.dumps(payload, ensure_ascii=False), request.headers.get("x-user", "")),
        )
    return get_creator(cid)


@app.post("/api/creators")
async def create_creator(request: Request) -> dict:
    payload = await request.json()
    handle = (payload.get("handle") or "").lstrip("@").strip()
    if not handle:
        raise HTTPException(400, "handle required")
    platform = payload.get("platform", "tiktok")
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute("SELECT 1 FROM creator WHERE platform = %s AND handle = %s", (platform, handle))
        if cur.fetchone():
            raise HTTPException(409, "creator already exists")
        fields = {k: v for k, v in payload.items() if k in CREATOR_EDITABLE | {"handle"}}
        fields["handle"] = handle
        fields.setdefault("platform", "tiktok")
        fields.setdefault("profile_url", f"https://www.tiktok.com/@{handle}")
        fields = encode_for_db(fields, CREATOR_JSON_COLS)
        cols = list(fields.keys())
        placeholders = ",".join(["%s"] * len(cols))
        cur.execute(
            f"INSERT INTO creator({','.join(cols)}) VALUES({placeholders}) RETURNING id",
            list(fields.values()),
        )
        cid = cur.fetchone()["id"]
        cur.execute(
            """
            INSERT INTO audit_log(table_name, record_id, action, changes, operator)
            VALUES('creator', %s, 'insert', %s, %s)
            """,
            (cid, json.dumps(fields, ensure_ascii=False), request.headers.get("x-user", "")),
        )
    return get_creator(cid)


@app.delete("/api/creators/{cid}")
def delete_creator(cid: int) -> dict:
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute("SELECT handle FROM creator WHERE id = %s", (cid,))
        r = cur.fetchone()
        if not r:
            raise HTTPException(404, "creator not found")
        cur.execute("DELETE FROM creator WHERE id = %s", (cid,))
        cur.execute(
            "INSERT INTO audit_log(table_name, record_id, action, operator) VALUES('creator', %s, 'delete', %s)",
            (cid, ""),
        )
        return {"ok": True, "deleted": r["handle"]}


# ---------- Outreach ----------
@app.get("/api/outreach")
def list_outreach(creator_id: int | None = None, limit: int = 200, offset: int = 0) -> dict:
    where = ["1=1"]
    args: list[Any] = []
    if creator_id is not None:
        where.append("o.creator_id = %s")
        args.append(creator_id)
    where_sql = " AND ".join(where)
    with get_pg_con() as con, con.cursor() as cur:
        cur.execute(
            f"""
            SELECT o.*, c.handle AS creator_handle
            FROM outreach o
            LEFT JOIN creator c ON c.id = o.creator_id
            WHERE {where_sql}
            ORDER BY o.event_date DESC NULLS LAST, o.id DESC
            LIMIT %s OFFSET %s
            """,
            args + [limit, offset],
        )
        rows = cur.fetchall()
        cur.execute(f"SELECT COUNT(*)::int AS total FROM outreach o WHERE {where_sql}", args)
        total = cur.fetchone()["total"]
        return {"total": total, "items": rows}


@app.post("/api/outreach")
async def create_outreach(request: Request) -> dict:
    payload = await request.json()
    fields = {k: v for k, v in payload.items() if k in OUTREACH_EDITABLE}
    if "creator_id" not in fields:
        raise HTTPException(400, "creator_id required")
    with get_pg_con() as con, con.cursor() as cur:
        cols = list(fields.keys())
        placeholders = ",".join(["%s"] * len(cols))
        cur.execute(
            f"INSERT INTO outreach({','.join(cols)}) VALUES({placeholders}) RETURNING id",
            list(fields.values()),
        )
        oid = cur.fetchone()["id"]
        if fields.get("status"):
            cur.execute(
                """
                UPDATE creator
                SET current_status = %s,
                    last_contact_date = COALESCE(%s, last_contact_date)
                WHERE id = %s
                """,
                (fields["status"], fields.get("event_date"), fields["creator_id"]),
            )
        cur.execute("SELECT * FROM outreach WHERE id = %s", (oid,))
        return cur.fetchone()


# ---------- Image serving ----------
def rel_path_to_url(rel: str) -> str:
    if rel.startswith("intern://"):
        return f"/intern/{rel[len('intern://'):]}"
    return f"/{rel}"


@app.get("/intern/{path:path}")
def serve_intern(path: str) -> FileResponse:
    decoded = unquote(path)
    target = (INTERN_A / decoded).resolve()
    try:
        target.relative_to(INTERN_A.resolve())
    except ValueError:
        raise HTTPException(403, "out of bounds")
    if not target.exists():
        raise HTTPException(404, str(target))
    return FileResponse(target)


# Mount static + assets
app.mount("/assets", StaticFiles(directory=str(ASSETS_DIR)), name="assets")
app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")


# ---------- Docs viewer (renders docs/*.md / *.docx / *.pdf / *.pptx / *.xlsx) ----------
import re as _re
import html as _html
from fastapi.responses import HTMLResponse, FileResponse
DOCS_DIR = ROOT / "docs"

# 允许从这些根目录服务文档
ALLOWED_DOC_ROOTS = [
    DOCS_DIR,
    Path(r"F:\实习生\C达人建联"),
    Path(r"F:\实习生\A社媒"),
]

_DOC_NAME_RE = _re.compile(
    r"^[\w一-鿿. \-（）()&]+\.(md|docx|pdf|pptx|xlsx|doc)$",
    _re.IGNORECASE
)

_DOC_HTML_SHELL = """<!doctype html>
<html lang="zh-CN"><head>
<meta charset="utf-8">
<title>{title}</title>
<style>
  body {{ max-width: 900px; margin: 30px auto; padding: 0 24px;
         font-family: -apple-system, "PingFang SC", "Microsoft YaHei", sans-serif;
         line-height: 1.7; color: #1f2329; }}
  h1, h2, h3, h4 {{ margin-top: 1.5em; line-height: 1.3; }}
  h1 {{ border-bottom: 2px solid #e5e6eb; padding-bottom: 8px; }}
  h2 {{ border-bottom: 1px solid #e5e6eb; padding-bottom: 6px; }}
  code {{ background: #f5f6f7; padding: 1px 5px; border-radius: 3px; font-size: 90%; }}
  pre {{ background: #1f2329; color: #fff; padding: 14px; border-radius: 6px;
        overflow-x: auto; font-size: 13px; line-height: 1.45; }}
  pre code {{ background: transparent; color: inherit; padding: 0; font-size: inherit; }}
  table {{ border-collapse: collapse; width: 100%; margin: 1em 0; font-size: 13px; }}
  th, td {{ border: 1px solid #e5e6eb; padding: 6px 10px; text-align: left; vertical-align: top; }}
  th {{ background: #fafbfc; font-weight: 600; }}
  a {{ color: #3370ff; text-decoration: none; }}
  a:hover {{ text-decoration: underline; }}
  blockquote {{ border-left: 3px solid #3370ff; padding: 6px 14px; color: #4e5969;
              margin: 1em 0; background: #eef2ff; border-radius: 0 4px 4px 0; }}
  ul, ol {{ padding-left: 28px; }}
  li {{ margin: 4px 0; }}
  hr {{ border: 0; border-top: 1px solid #e5e6eb; margin: 2em 0; }}
  img {{ max-width: 100%; }}
  .nav-back {{ position: fixed; top: 16px; right: 20px; padding: 6px 12px;
              background: #3370ff; color: #fff; border-radius: 4px;
              font-size: 12px; z-index: 99; box-shadow: 0 2px 6px rgba(0,0,0,0.1); }}
  .file-badge {{ display: inline-block; padding: 3px 9px; border-radius: 11px;
                background: #eef2ff; color: #3370ff; font-size: 11px; margin-left: 8px;
                vertical-align: middle; }}
  .source-note {{ background: #fafbfc; border: 1px solid #e5e6eb; border-radius: 6px;
                 padding: 8px 12px; margin-bottom: 14px; font-size: 12px; color: #4e5969; }}
  .download-link {{ float: right; padding: 4px 10px; background: #3370ff; color: #fff;
                    border-radius: 4px; font-size: 11px; }}
  .slide {{ border: 1px solid #e5e6eb; border-radius: 6px; padding: 12px 16px;
           margin: 12px 0; background: #fafbfc; }}
</style></head>
<body>
<a class="nav-back" href="/">← 返回主界面</a>
{body}
</body></html>"""


def _find_doc(name: str) -> Path | None:
    """Search whitelisted roots for the requested file name."""
    for root in ALLOWED_DOC_ROOTS:
        if not root.exists():
            continue
        candidate = root / name
        try:
            resolved = candidate.resolve()
            resolved.relative_to(root.resolve())
        except (ValueError, OSError):
            continue
        if resolved.is_file():
            return resolved
    return None


def _render_md(fp: Path, name: str) -> HTMLResponse:
    try:
        import markdown as _md
        body = _md.markdown(fp.read_text(encoding="utf-8"),
                            extensions=["tables", "fenced_code", "sane_lists", "toc"])
    except ImportError:
        body = "<pre>" + _html.escape(fp.read_text(encoding="utf-8")) + "</pre>"
    return HTMLResponse(_DOC_HTML_SHELL.format(title=name, body=body))


def _render_docx(fp: Path, name: str) -> HTMLResponse:
    try:
        from docx import Document
    except ImportError:
        return FileResponse(fp, filename=fp.name)
    doc = Document(str(fp))
    parts = [f'<h1>{_html.escape(name)} <span class="file-badge">.docx</span></h1>',
             '<div class="source-note">📄 Word 文档 — 已抽取文字内容（保留段落 / 标题 / 表格；样式简化）。'
             f'<a class="download-link" href="/docs/raw/{name}" download>下载原文件</a></div>']
    # 段落
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if "heading 1" in style:
            parts.append(f"<h2>{_html.escape(text)}</h2>")
        elif "heading 2" in style:
            parts.append(f"<h3>{_html.escape(text)}</h3>")
        elif "heading" in style:
            parts.append(f"<h4>{_html.escape(text)}</h4>")
        else:
            parts.append(f"<p>{_html.escape(text)}</p>")
    # 表格
    for ti, table in enumerate(doc.tables, 1):
        parts.append(f'<h3>表 {ti}</h3><table>')
        for ri, row in enumerate(table.rows):
            tag = "th" if ri == 0 else "td"
            cells = "".join(
                f"<{tag}>{_html.escape(c.text.strip())}</{tag}>" for c in row.cells
            )
            parts.append(f"<tr>{cells}</tr>")
        parts.append("</table>")
    return HTMLResponse(_DOC_HTML_SHELL.format(title=name, body="\n".join(parts)))


def _render_pptx(fp: Path, name: str) -> HTMLResponse:
    try:
        from pptx import Presentation
    except ImportError:
        return FileResponse(fp, filename=fp.name)
    prs = Presentation(str(fp))
    parts = [f'<h1>{_html.escape(name)} <span class="file-badge">.pptx</span></h1>',
             '<div class="source-note">📊 PowerPoint — 已抽取每张幻灯片的文字内容。'
             f'<a class="download-link" href="/docs/raw/{name}" download>下载原文件</a></div>']
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'<div class="slide"><h3>Slide {i}</h3>')
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text.strip():
                continue
            for line in shape.text.splitlines():
                if line.strip():
                    parts.append(f"<p>{_html.escape(line.strip())}</p>")
        parts.append("</div>")
    return HTMLResponse(_DOC_HTML_SHELL.format(title=name, body="\n".join(parts)))


def _render_xlsx(fp: Path, name: str) -> HTMLResponse:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return FileResponse(fp, filename=fp.name)
    wb = load_workbook(str(fp), read_only=True, data_only=True)
    parts = [f'<h1>{_html.escape(name)} <span class="file-badge">.xlsx</span></h1>',
             '<div class="source-note">📈 Excel — 已渲染每个 sheet 前 200 行。'
             f'<a class="download-link" href="/docs/raw/{name}" download>下载原文件</a></div>']
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"<h3>Sheet: {_html.escape(sheet_name)}</h3>")
        parts.append("<table>")
        for ri, row in enumerate(ws.iter_rows(max_row=200, values_only=True)):
            tag = "th" if ri == 0 else "td"
            cells = "".join(
                f"<{tag}>{_html.escape(str(c) if c is not None else '')}</{tag}>" for c in row
            )
            parts.append(f"<tr>{cells}</tr>")
        parts.append("</table>")
    wb.close()
    return HTMLResponse(_DOC_HTML_SHELL.format(title=name, body="\n".join(parts)))


@app.get("/docs/raw/{name}")
def serve_doc_raw(name: str) -> FileResponse:
    """Download the original file (used by the 'download original' link in viewers)."""
    if not _DOC_NAME_RE.match(name):
        raise HTTPException(400, "invalid doc name")
    fp = _find_doc(name)
    if not fp:
        raise HTTPException(404, f"doc '{name}' not found")
    return FileResponse(fp, filename=fp.name)


@app.get("/docs/{name}")
def serve_doc(name: str):
    """Render docs/*.md, docx/pdf/pptx/xlsx in browser; falls back to download for unknown types.

    Searches across whitelisted roots: docs/, F:\\实习生\\C达人建联, F:\\实习生\\A社媒
    """
    if not _DOC_NAME_RE.match(name):
        raise HTTPException(400, f"invalid doc name '{name}' (must end in .md/.docx/.pdf/.pptx/.xlsx/.doc)")
    fp = _find_doc(name)
    if not fp:
        raise HTTPException(
            404,
            f"'{name}' not found in any allowed location. "
            f"Searched: docs/ + 实习生/C达人建联 + 实习生/A社媒"
        )
    ext = fp.suffix.lower()
    if ext == ".md":
        return _render_md(fp, name)
    if ext == ".docx":
        return _render_docx(fp, name)
    if ext == ".pptx":
        return _render_pptx(fp, name)
    if ext == ".xlsx":
        return _render_xlsx(fp, name)
    if ext == ".pdf":
        return FileResponse(fp, media_type="application/pdf", filename=fp.name)
    # .doc legacy etc — just serve as download
    return FileResponse(fp, filename=fp.name)
